from uagents import Agent, Context, Model, Protocol, Bureau
from pptx import Presentation
import io
import base64
import os
import tempfile
import threading
import mimetypes
from openai import OpenAI
from dotenv import load_dotenv
from pathlib import Path

# Load .env file
load_dotenv()

api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)
class Request(Model):
    file_data: str
    file_name: str
    file_type: str = "unknown"

class Response(Model):
    message: str
    success: bool

class ProcessedContent(Model):
    slide_content: list
    slide_count: int


# ####################################
# Agents Initialization
# ####################################

# Root agent 
# This agent serves as the main entry point for the application.
root_agent = Agent(name="root_agent")

# File classifier agent
# This agent serves as a classifier for the uploaded files. 
# It will classify the file type and send it to the appropriate agent for further processing.
file_classifier_agent = Agent(name="file classifier agent")

# PPT processing agent
# This agent will prep the ppt file beofre sending it to the LLM for further processing.
ppt_processing_agent = Agent(name="ppt processing agent")

# Summary agent
# This agent will summarize the processed content.
summary_agent = Agent(name="summary agent")

# ####################################
# Agents Startup 
# ####################################

@root_agent.on_event("startup")
async def startup_event(ctx: Context):
    """
    Log a message when the root agent starts up.

    """
    ctx.logger.info(f"Root agent {root_agent.name} started with address {root_agent.address}")

@file_classifier_agent.on_event("startup") 
async def classifier_startup(ctx: Context):
    ctx.logger.info(f"File classifier agent {file_classifier_agent.name} started with address {file_classifier_agent.address}")

@ppt_processing_agent.on_event("startup")
async def ppt_processing_startup(ctx: Context):
    ctx.logger.info(f"PPT processing agent {ppt_processing_agent.name} started with address {ppt_processing_agent.address}")

@summary_agent.on_event("startup")
async def summary_startup(ctx: Context):
    ctx.logger.info(f"Summary agent {summary_agent.name} started with address {summary_agent.address}")
 

# ####################################
# Agents Event Handlers
# ####################################
@root_agent.on_rest_post("/rest/post", Request, Response)
async def handle_file_upload(ctx: Context, req: Request):
    """
    Starting point for the file processing workflow
    """
    try:
        ctx.logger.info(f"Received file upload request: {req.file_name}")

        # Forwared to file classifier agent
        response = await ctx.send(file_classifier_agent.address, req)
    except Exception as e:
        ctx.logger.error(f"Error processing file upload: {e}")
        return Response(message=str(e), success=False)



@file_classifier_agent.on_message(model=Request)
async def classify_file(ctx: Context, sender: str, req: Request):
    """
    Classify the uploaded file and forward it to the appropriate agent.
    """

    ctx.logger.info(f"Classifying file: {req.file_name}")

    # Check the file type and forward to the appropriate agent
    if is_ppt_file(req.file_name):
        ctx.logger.info(f"Detected file {req.file_name} is a PowerPoint file, forwarding to PPT processing agent.")
        response = await ctx.send(ppt_processing_agent.address, req)

@ppt_processing_agent.on_message(model=Request)
async def process_ppt_file(ctx: Context, sender: str, req: Request):
    """
    Process the PowerPoint file and prepare it for further analysis.
    """
    ctx.logger.info(f"Processing PowerPoint file: {req.file_name}")

 
    # Decode the base64 file data
    file_content = base64.b64decode(req.file_data)
    
    extracted_contents = extract_text_from_ppt(file_content)
    processed_content = ProcessedContent(
        slide_content=extracted_contents['slide_content'],
        slide_count=extracted_contents['slide_count']
    )
    await ctx.send(summary_agent.address, processed_content)

@summary_agent.on_message(model=ProcessedContent)
async def summarize_content(ctx: Context, sender: str, content: ProcessedContent):
    """
    Summarize the processed content using Open AI. 
    """
    ctx.logger.info(f"Summarizing content with {len(content.slide_content)} slides.")

    # Batch summarize the slide content
    summaries = batch_summarize(content.slide_content, content.slide_count)

    script = generate_script(summaries)
    if len(script) != 0:
        # Use pathlib to write the script content to a file
        output_dir = Path(__file__).parent / "script_output"
        output_dir.mkdir(exist_ok=True)
        file_path = output_dir / "script.txt"
        file_path.write_text(script)



    

# ####################################
# Helper functions
# ####################################
def is_ppt_file(file_name: str) -> bool:
    """
    Check if the file is a PowerPoint file based on its extension.
    """
    mime_type, _ = mimetypes.guess_type(file_name)
    return mime_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]

def extract_text_from_ppt(file_content: bytes) -> object:
    """
    Extract text from a PowerPoint file.
    This is a placeholder function. You can implement actual PPT processing here.
    """
    # Load the presentation from the byte content
    prs = Presentation(io.BytesIO(file_content))

    extracted_data = {
        'slide_content': [],
        'slide_count': len(prs.slides),
        
    }

   

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'content': '',
            'notes': '',
            'all_text': ''
        }
        slide_texts = []

        # Collect all text in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_texts.append(text)
                
                # Try to identify title (usually first text or largest)
                if not slide_data['title'] and len(text) < 100:
                    slide_data['title'] = text
                else:
                    slide_data['content'] += text + '\n'

        # Collect text in speaker note
        if slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data['notes'] = notes_text
                slide_texts.append(f"[Speaker Notes: {notes_text}]")
        
        slide_data['all_text'] = '\n'.join(slide_texts)
        extracted_data['slide_content'].append(slide_data)
    return extracted_data

def batch_summarize(slide_content, total_slide, batch_size=5, model="gpt-5-nano"): 
    """
    Batch summarize the slide content.
    
    slide_content: [
        {
            'slide_number': slide_num,
            'title': '',
            'content': '',
            'notes': '',
            'all_text': ''
        },
        ...
    ]
    """
    summaries = []

    for i in range(0, total_slide, batch_size):
        batch = slide_content[i:i + batch_size]
        
        # Create prompt
        slides_text = "\n\n".join([f"Slide {slide['slide_number']}: {slide['content']}" for slide in batch])
        
        prompt = f"""
        Summarize the following slides briefly.
        Use one bullet point per slide. Keep text minimal but clear
        
        Response should be in the following format:
        - Slide 1: Summary of slide 1 content
        - Slide 2: Summary of slide 2 content
        ...

        {slides_text}
        """

        # Use GPT-5 to summarize the batch

        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You are a summarizer. Always produce the shortest possible summaries."},
                {"role": "user", "content": prompt}
            ]
        )
        batch_summary = response.choices[0].message.content.strip()
        # Clean & split into usable summaries
        lines = batch_summary.split("\n")
        for line in lines:
            line = line.strip("-• ").strip()
            if line:
                summaries.append(line)
    
    return summaries



def generate_script(summaries, model='gpt-5-nano'):
    """
    Generate a script from the summaries using GPT-5.
    """
    summaries_text = "\n".join([f"- {s}" for s in summaries])

    prompt = f"""
    You are a presentation script writer creating a spoken script for a presenter.

    Instructions:
    - Write natural, conversational scripts as if speaking to an audience
    - Start each script section with "Slide(s) X:" 
    - Reference the slide number(s) naturally in the script (e.g., "As we can see in Slide 3...")
    - Keep each script segment 2-4 sentences long
    - Use engaging, clear language appropriate for the audience
    - Group related slides together when it makes sense

    Slide Content:
    {summaries_text}

    Required Output Format:
    Slide(s) 1:
    [Script text mentioning Slide 1]

    Slide(s) 2:
    [Script text for slide 2]

    Slide(s) 3:
    [Script text for slide 3]

    Slide(s) 4:
    [Script text for slide 4]

    Continue this pattern for all slides.
    """

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": "You are a script writer for presentations."},
            {"role": "user", "content": prompt}
        ]
    )

    return response.choices[0].message.content.strip()

    
        
       





# ####################################
# Bureau Setup
# ####################################
if __name__ == "__main__":
    bureau = Bureau(port=8080)
    bureau.add(root_agent)
    bureau.add(file_classifier_agent)
    bureau.add(ppt_processing_agent)
    bureau.add(summary_agent)
    bureau.run()